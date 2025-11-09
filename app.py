import os
import re
import uuid
from pathlib import Path
from typing import List
from flask import Flask, render_template, request, send_from_directory, redirect, url_for, flash
import subprocess
import cv2
import json
from PIL import Image
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches

# ------------------- CONFIG -------------------
BLUR_LAPLACIAN_THRESH = 5.0
SAFE_FORMATS = ["best", "bestvideo+bestaudio", "bestvideo[ext=mp4]+bestaudio[ext=m4a]/mp4"]
WORKDIR = Path("runs")
WORKDIR.mkdir(exist_ok=True)
# ----------------------------------------------

app = Flask(__name__)
app.secret_key = os.environ.get("SECRET_KEY", "dev-secret-key")


def download_video(url: str, out_dir: Path) -> Path:
    out_dir.mkdir(parents=True, exist_ok=True)
    out_tmpl = str(out_dir / "%(title)s.%(ext)s")
    for fmt in SAFE_FORMATS:
        try:
            subprocess.run(["python3", "-m", "yt_dlp", "-f", fmt, "-o", out_tmpl, url], check=True)
            vids = sorted(out_dir.glob("*.mp4"), key=lambda p: p.stat().st_size, reverse=True)
            if vids:
                return vids[0]
        except subprocess.CalledProcessError:
            continue
    raise RuntimeError("All formats failed to download.")


def parse_time_str(tstr: str) -> float:
    parts = [float(p) for p in tstr.split(":")]
    if len(parts) == 1:
        return parts[0]
    elif len(parts) == 2:
        m, s = parts
        return m * 60 + s
    elif len(parts) == 3:
        h, m, s = parts
        return h * 3600 + m * 60 + s
    else:
        raise ValueError(f"Invalid time format: {tstr}")


def read_frame_at(video_path: Path, time_s: float):
    cap = cv2.VideoCapture(str(video_path))
    if not cap.isOpened():
        return None
    fps = cap.get(cv2.CAP_PROP_FPS) or 30
    total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
    duration = total_frames / fps if fps else 0
    if time_s > duration and duration > 0:
        time_s = max(0, duration - 0.1)
    frame_idx = max(0, min(int(time_s * fps), total_frames - 1))
    cap.set(cv2.CAP_PROP_POS_FRAMES, frame_idx)
    ok, frame = cap.read()
    if not ok:
        for i in range(1, 5):
            cap.set(cv2.CAP_PROP_POS_FRAMES, frame_idx - i)
            ok, frame = cap.read()
            if ok:
                break
    cap.release()
    return frame if ok else None


def blur_score_laplacian(gray):
    return cv2.Laplacian(gray, cv2.CV_64F).var()


def save_image(img_bgr, out_path: Path):
    out_path.parent.mkdir(parents=True, exist_ok=True)
    Image.fromarray(cv2.cvtColor(img_bgr, cv2.COLOR_BGR2RGB)).save(out_path, quality=95)


def make_pdf(images: List[Path], output_pdf: Path):
    pdf = FPDF()
    for img_path in images:
        img = Image.open(img_path)
        w, h = img.size
        orientation = 'L' if w > h else 'P'
        pdf.add_page(orientation=orientation)
        pdf.image(str(img_path), x=10, y=10, w=pdf.w - 20)
    pdf.output(str(output_pdf))


def make_ppt(images: List[Path], output_ppt: Path):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    for img_path in images:
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(str(img_path), Inches(0), Inches(0),
                                 width=prs.slide_width, height=prs.slide_height)
    prs.save(str(output_ppt))


@app.route("/", methods=["GET"])
def index():
    return render_template("index.html")


@app.route("/process", methods=["POST"])
def process():
    url = request.form.get("url", "").strip()
    times_raw = request.form.get("times", "").strip()
    want_pdf = request.form.get("pdf") == "on"
    want_ppt = request.form.get("ppt") == "on"
    if not url or not times_raw:
        flash("Please provide both the YouTube URL and the timestamps.", "error")
        return redirect(url_for("index"))

    sess_id = uuid.uuid4().hex[:8]
    sess_dir = WORKDIR / f"run_{sess_id}"
    frames_dir = sess_dir / "images"
    sess_dir.mkdir(parents=True, exist_ok=True)

    try:
        video_path = Path(url) if Path(url).exists() else download_video(url, sess_dir)
    except Exception as e:
        flash(f"Video download failed: {e}", "error")
        return redirect(url_for("index"))

    # Extract frames
    saved_images: List[Path] = []
    try:
        # parse times separated by comma/space/newline
        parts = re.split(r"[,\\s]+", times_raw)
        custom_times = [parse_time_str(x.strip()) for x in parts if x.strip()]
    except Exception as e:
        flash(f"Invalid time format: {e}", "error")
        return redirect(url_for("index"))

    for t in custom_times:
        frame = read_frame_at(video_path, t)
        if frame is None:
            continue
        gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
        bscore = blur_score_laplacian(gray)
        pass
        out_path = frames_dir / f"frame_{int(t*1000)}ms.jpg"
        save_image(frame, out_path)
        saved_images.append(out_path)

    # Delete downloaded video if it was remote
    if not Path(url).exists() and video_path.exists():
        try:
            video_path.unlink(missing_ok=True)
        except Exception:
            pass

    pdf_path = None
    ppt_path = None
    if saved_images:
        if want_pdf:
            pdf_path = sess_dir / "frames.pdf"
            make_pdf(saved_images, pdf_path)
        if want_ppt:
            ppt_path = sess_dir / "frames.pptx"
            make_ppt(saved_images, ppt_path)

    # Build a little manifest
    manifest = {
        "session": sess_id,
        "images": [str(p.name) for p in saved_images],
        "pdf": str(pdf_path.name) if pdf_path else None,
        "ppt": str(ppt_path.name) if ppt_path else None
    }
    (sess_dir / "manifest.json").write_text(json.dumps(manifest, indent=2))

    return redirect(url_for("results", sess_id=sess_id))


@app.route("/results/<sess_id>")
def results(sess_id):
    sess_dir = WORKDIR / f"run_{sess_id}"
    if not sess_dir.exists():
        return "Session not found", 404
    manifest_path = sess_dir / "manifest.json"
    manifest = json.loads(manifest_path.read_text()) if manifest_path.exists() else {}
    return render_template("results.html", sess_id=sess_id, manifest=manifest)


@app.route("/download/<sess_id>/<path:filename>")
def download(sess_id, filename):
    sess_dir = WORKDIR / f"run_{sess_id}"
    return send_from_directory(sess_dir, filename, as_attachment=True)


@app.route("/image/<sess_id>/<path:filename>")
def serve_image(sess_id, filename):
    frames_dir = WORKDIR / f"run_{sess_id}" / "images"
    return send_from_directory(frames_dir, filename)

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=int(os.environ.get("PORT", 7860)))
